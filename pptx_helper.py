import os
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont

# Intentar importar win32com para conversión nativa usando PowerPoint en Windows
try:
    import win32com.client
    HAS_WIN32COM = True
except ImportError:
    HAS_WIN32COM = False


def convert_pptx_to_images(pptx_path, output_dir):
    """
    Convierte un archivo PPTX/PPT en una serie de imágenes PNG.
    Retorna una lista de rutas de archivos de imágenes (vacía si falla todo).
    """
    os.makedirs(output_dir, exist_ok=True)
    images = []

    # Método 1: Conversión nativa con MS PowerPoint instalado (mejor calidad)
    if HAS_WIN32COM:
        images = _try_win32com(pptx_path, output_dir)
        if images:
            return images

    # Método 2: Fallback con python-pptx + Pillow (solo funciona con .pptx, NO .ppt)
    if pptx_path.lower().endswith(".pptx"):
        images = _try_pptx_fallback(pptx_path, output_dir)
        if images:
            return images
        else:
            print("ADVERTENCIA: El fallback de texto no generó diapositivas. "
                  "Verifica que el archivo .pptx no esté corrupto o protegido.")
    else:
        print(f"ADVERTENCIA: Los archivos .ppt (formato antiguo) no se pueden leer sin "
              f"Microsoft PowerPoint instalado. Convierte '{os.path.basename(pptx_path)}' "
              f"a .pptx o impórtalo como PDF en su lugar.")

    return images


def _try_win32com(pptx_path, output_dir):
    powerpoint = None
    presentation = None
    images = []
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        abs_pptx_path = os.path.abspath(pptx_path)
        abs_output_dir = os.path.abspath(output_dir)

        presentation = powerpoint.Presentations.Open(abs_pptx_path, WithWindow=False)
        presentation.SaveAs(abs_output_dir, 17)  # 17 = formato PNG

        presentation.Close()
        powerpoint.Quit()

        files = os.listdir(output_dir)
        png_files = [f for f in files if f.lower().endswith(".png")]

        def extract_number(filename):
            num_str = "".join([c for c in filename if c.isdigit()])
            return int(num_str) if num_str else 0

        png_files.sort(key=extract_number)
        images = [os.path.join(output_dir, png) for png in png_files]

        if images:
            print(f"PowerPoint convertido vía win32com: {len(images)} diapositivas.")
        return images

    except Exception as e:
        print(f"win32com no disponible o falló ({e}). Probando fallback...")
        try:
            if presentation:
                presentation.Close()
        except Exception:
            pass
        try:
            if powerpoint:
                powerpoint.Quit()
        except Exception:
            pass
        return []


def _try_pptx_fallback(pptx_path, output_dir):
    """Extrae texto de cada diapositiva y lo dibuja sobre un fondo oscuro limpio."""
    images = []
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"No se pudo abrir el .pptx con python-pptx: {e}")
        return []

    font = _load_font(48)

    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        text_content = "\n".join(slide_text) if slide_text else f"Diapositiva {i + 1}"

        width, height = 1920, 1080
        img = Image.new("RGB", (width, height), color="#121214")
        draw = ImageDraw.Draw(img)

        lines = [l.strip() for l in text_content.split("\n") if l.strip()]
        if not lines:
            lines = [f"Diapositiva {i + 1}"]

        # Calcular bloque centrado verticalmente
        line_heights = []
        for line in lines[:8]:
            bbox = draw.textbbox((0, 0), line, font=font)
            line_heights.append(bbox[3] - bbox[1])
        total_h = sum(line_heights) + (len(lines[:8]) - 1) * 30
        y_offset = (height - total_h) // 2

        for idx, line in enumerate(lines[:8]):
            bbox = draw.textbbox((0, 0), line, font=font)
            text_w = bbox[2] - bbox[0]
            text_h = bbox[3] - bbox[1]
            x = (width - text_w) // 2
            draw.text((x, y_offset), line, font=font, fill="#ffffff")
            y_offset += text_h + 30

        img_path = os.path.join(output_dir, f"slide_{i + 1}.png")
        img.save(img_path)
        images.append(img_path)

    if images:
        print(f"PowerPoint procesado con python-pptx (fallback de texto): {len(images)} diapositivas.")
    return images


def _load_font(size):
    """Intenta cargar una fuente legible; si no encuentra ninguna, usa la default de Pillow."""
    candidates = [
        "arial.ttf",
        "C:/Windows/Fonts/arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        "/System/Library/Fonts/Supplemental/Arial.ttf",
    ]
    for path in candidates:
        try:
            return ImageFont.truetype(path, size)
        except (IOError, OSError):
            continue
    return ImageFont.load_default()